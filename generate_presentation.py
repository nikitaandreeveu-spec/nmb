#!/usr/bin/env python3
"""
Generate Anthropic Business Model PowerPoint Presentation
Creates a 11-slide presentation with Anthropic aesthetic (Beige, Safety Green, Orange)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# ─── Anthropic Color Palette ──────────────────────────────────────
BEIGE = RGBColor(245, 242, 237)          # #F5F2ED
CHARCOAL = RGBColor(26, 26, 26)          # #1A1A1A
SAFETY_GREEN = RGBColor(74, 124, 89)     # #4A7C59
ACCENT_ORANGE = RGBColor(196, 98, 45)    # #C4622D
SURFACE = RGBColor(255, 255, 255)        # #FFFFFF
BORDER = RGBColor(224, 221, 216)         # #E0DDD8
LIGHT_GREEN = RGBColor(61, 102, 73)      # #3d6649

# ─── Create Presentation ──────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BEIGE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.margin_bottom = Inches(0.2)
    
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = CHARCOAL
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.LEFT
    
    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = SAFETY_GREEN
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.LEFT
    
    return slide

def add_content_slide(prs, title, content_lines, highlight_items=None):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BEIGE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = CHARCOAL
    p.font.name = 'Calibri'
    
    # Add top border line
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.3), Inches(9.5), Inches(1.3))
    line.line.color.rgb = SAFETY_GREEN
    line.line.width = Pt(3)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(8.6), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = CHARCOAL
        p.font.name = 'Calibri'
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.level = 0
        
        # Highlight certain items
        if highlight_items and any(item in line for item in highlight_items):
            p.font.bold = True
            p.font.color.rgb = SAFETY_GREEN
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BEIGE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = CHARCOAL
    p.font.name = 'Calibri'
    
    # Add top border line
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.3), Inches(9.5), Inches(1.3))
    line.line.color.rgb = SAFETY_GREEN
    line.line.width = Pt(3)
    
    # Add table
    rows_count = len(rows) + 1
    cols_count = len(headers)
    left = Inches(0.7)
    top = Inches(1.7)
    width = Inches(8.6)
    height = Inches(5.2)
    
    table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table
    
    # Set column widths (convert to EMUs - English Metric Units)
    col_width = int(width / cols_count)
    for col_idx in range(cols_count):
        table_shape.columns[col_idx].width = col_width
    
    # Header row
    for col_idx, header in enumerate(headers):
        cell = table_shape.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = SAFETY_GREEN
        
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = SURFACE
            paragraph.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row_data in enumerate(rows, 1):
        for col_idx, cell_data in enumerate(row_data):
            cell = table_shape.cell(row_idx, col_idx)
            cell.text = str(cell_data)
            cell.fill.solid()
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(245, 247, 242)  # Light green tint
            else:
                cell.fill.fore_color.rgb = SURFACE
            
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = CHARCOAL
                paragraph.alignment = PP_ALIGN.CENTER
    
    return slide

# ─── SLIDE 1: Title/Cover ────────────────────────────────────────
slide = add_title_slide(prs, "Anthropic PBC", "The Safety-Frontier Business Model (2026)")
footer = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.4))
p = footer.text_frame.paragraphs[0]
p.text = "Building the responsible, outcome-driven enterprise AI infrastructure of tomorrow"
p.font.size = Pt(14)
p.font.italic = True
p.font.color.rgb = SAFETY_GREEN

# ─── SLIDE 2: Business Model Canvas Overview ────────────────────
slide = add_content_slide(prs, "Business Model Canvas: Value System Overview",
    [
        "• Economic Value: Claude solves real professional problems in coding, research, strategy",
        "• Social Value: Public Benefit Corporation structure signals safety-first commitment",
        "• Systemic Value: Responsible Scaling Policy manages risk mitigation and resilience",
        "• Brand Value: Ad-free 'space to think' protects integrity of user-model interaction",
        "",
        "Key Insight: Anthropic's value proposition spans user performance, societal responsibility, systemic risk management, and brand trust—not just AI software delivery."
    ],
    highlight_items=["Economic Value", "Social Value", "Systemic Value", "Brand Value"]
)

# ─── SLIDE 3: Strategic Innovation Pathways ────────────────────
slide = add_content_slide(prs, "Strategic Innovation Pathways",
    [
        "Pathway 1: Digital Provenance & Audit Trail",
        "  → Financial auditing, regulatory compliance, accountability tracking",
        "",
        "Pathway 2: Governance-as-a-Service (GaaS)",
        "  → Position Anthropic as 'Safety OS' for entire AI ecosystem",
        "",
        "Pathway 3: Outcome-Based 'Agentic Workforce' ⭐ MOST VIABLE",
        "  → Pay for measurable results, not compute access",
        "  → Addresses 2026 demand for hard ROI and trust",
        "",
        "Strategic Focus: Combine efficiency excellence with outcome-based pricing to create uncontested market position."
    ],
    highlight_items=["MOST VIABLE", "Outcome-Based"]
)

# ─── SLIDE 4: PESTLE Analysis Summary ────────────────────────────
slide = add_content_slide(prs, "PESTLE Analysis: Strategic Context",
    [
        "Political: AI sovereignty demands & geopolitical alignment create regulatory moat",
        "Economic: Monetization shift → outcome-based pricing; $5-30B capital needs",
        "Social: Trust gap, labor concerns, public-benefit mission as differentiator",
        "Technological: Agentic AI frontier, vertical product success (Claude Code: 54% market share)",
        "Legal: IP/copyright litigation, compliance frameworks becoming competitive moat",
        "Environmental: Energy consumption, sustainability commitments, ESG differentiation"
    ],
    highlight_items=["regulatory moat", "outcome-based pricing", "54% market share", "competitive moat"]
)

# ─── SLIDE 5: PESTLE-BMC Impact Matrix ────────────────────────
slide = add_table_slide(prs, "PESTLE-BMC Impact Matrix: Critical Zones",
    ["Business Model Component", "Geopolitics", "Compute Power", "Compliance"],
    [
        ["Key Partnerships", "🔴 CRITICAL", "🔴 CRITICAL", "🟠 HIGH"],
        ["Value Propositions", "🔴 CRITICAL", "🔴 CRITICAL", "🟠 HIGH"],
        ["Key Resources", "🔴 CRITICAL", "🔴 CRITICAL", "🟠 HIGH"],
        ["Cost Structure", "🔴 CRITICAL", "🔴 CRITICAL", "🟠 HIGH"],
        ["Revenue Streams", "🔴 CRITICAL", "🔴 CRITICAL", "🟠 HIGH"]
    ]
)

# ─── SLIDE 6: Future Cone Scenarios ──────────────────────────────
slide = add_content_slide(prs, "Strategic Foresight: Three Future Scenarios",
    [
        "Probable Future (2026): 'The Digital Coworker'",
        "  → AI evolves into partner automating multi-step processes",
        "",
        "Plausible Future (2028): 'The Trust Crisis'",
        "  → Agentic loops trigger accountability crisis; 1 in 3 consumers trust AI",
        "",
        "Preferable Future (Our Goal): 'The Outcome Economy'",
        "  → AI as reliable 'Invisible Workforce'; pay only for measurable results",
        "",
        "Strategic Path: Drive toward Preferable Future through Constitutional AI, compliance products, and outcome-based pricing."
    ],
    highlight_items=["Digital Coworker", "Trust Crisis", "Outcome Economy", "Preferable Future"]
)

# ─── SLIDE 7: Executive Summary - The Pitch ──────────────────────
slide = add_content_slide(prs, "Executive Summary: Business Model Transformation",
    [
        "$50B+ Incremental Value by 2028 through Three Strategic Shifts:",
        "",
        "1️⃣  Revenue Model: Seat-Based → Outcome-Based (+$8-12B revenue)",
        "2️⃣  Technology: Performance-First → Efficiency-First (40-50% cost reduction)",
        "3️⃣  Customer Engagement: Direct → Enterprise Partnerships (+$15-20B TAM)",
        "",
        "Financial Projections: $9B (2025) → $25-30B (2028) • 38-42% CAGR",
        "Gross Margin: 45-50% (2025) → 60-75% (2028)"
    ],
    highlight_items=["$50B+", "Three Strategic Shifts", "$25-30B", "60-75%"]
)

# ─── SLIDE 8: Why Anthropic Wins ──────────────────────────────
slide = add_content_slide(prs, "Competitive Differentiation: Why Anthropic Wins",
    [
        "✓ Trust & Compliance Moat: PBC structure + Constitutional AI",
        "  → Regulatory compliance becomes profitable product ($2-3B potential)",
        "",
        "✓ Efficiency Excellence: 3-4x improvements per generation",
        "  → Green AI unlocks sustainable margin expansion",
        "",
        "✓ Vertical Specialization: Claude Code at 54% enterprise market",
        "  → Domain-specific agents command 3-5x pricing premiums",
        "",
        "✓ Sovereign AI Leadership: Claude Gov + allied nation partnerships",
        "  → Government revenue $3-5B by 2028"
    ],
    highlight_items=["Trust & Compliance", "Efficiency Excellence", "Vertical Specialization", "Sovereign AI"]
)

# ─── SLIDE 9: Implementation Timeline Phase 1 ────────────────────
slide = add_content_slide(prs, "Phase 1: Foundation & Pilot (2026)",
    [
        "Q1: Pricing framework, partnerships, sovereign cloud setup",
        "",
        "Q2: 10-15 outcome pilots, Accenture deal ($5-10M), Claude Gov expansion",
        "",
        "Q3: 50-100 customers migrated, 3 verticals launched, compute partnerships",
        "",
        "Q4: $15M+ ARR achieved, $50M+ vertical run rate, org scaling complete",
        "",
        "Target: Foundation built for Phase 2 acceleration"
    ]
)

# ─── SLIDE 10: Implementation Timeline Phase 2-3 ──────────────────
slide = add_table_slide(prs, "Phase 2-3 Goals & Financial Trajectory",
    ["Metric", "Phase 1 (2026)", "Phase 2 (2027)", "Phase 3 (2028)"],
    [
        ["Annual Revenue", "$12-15B", "$18-22B", "$25-30B"],
        ["Gross Margin", "48-52%", "55-60%", "60-75%"],
        ["Outcome Contracts", "150-200", "300-400", "600-700"],
        ["Gov Revenue", "$200M", "$1B", "$2-3B"]
    ]
)

# ─── SLIDE 11: Digital Infrastructure & Risks ─────────────────
slide = add_content_slide(prs, "Digital Infrastructure & Implementation Risks",
    [
        "6 Core Platforms: EOMP (outcomes), VADP (verticals), RCSAP (compliance),",
        "GAOP (green AI), GSCOC (government), CPEP (partner enablement)",
        "",
        "🔴 CRITICAL Risks: Compute supply, geopolitics, organizational scaling",
        "🟠 HIGH Risks: Outcome pricing adoption, vertical development velocity",
        "",
        "3-Year Investment: $324-427M (infrastructure, talent, partnerships, R&D)",
        "",
        "Success Factors: Partner execution, regulatory clarity, compute supply, talent velocity, customer adoption"
    ],
    highlight_items=["Core Platforms", "CRITICAL Risks", "HIGH Risks", "$324-427M"]
)

# ─── Save Presentation ──────────────────────────────────────────
output_path = '/workspaces/nmb/Anthropic_Business_Model_2026.pptx'
prs.save(output_path)
print(f"✅ PowerPoint presentation created: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
